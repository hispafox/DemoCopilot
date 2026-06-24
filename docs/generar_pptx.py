"""Genera analisis-diseño.pptx con python-pptx."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree

OUTPUT = os.path.join(os.path.dirname(__file__), "analisis-diseño.pptx")

# ── Paleta "Midnight Executive" adaptada a .NET ───────────────────────────────
AZUL_OSC  = RGBColor(0x1E, 0x27, 0x61)   # fondo portada / slides dark
AZUL_MED  = RGBColor(0x2E, 0x75, 0xB6)   # cabeceras tablas, acentos
AZUL_CLAR = RGBColor(0xCA, 0xDC, 0xFC)   # filas pares tablas, pills
BLANCO    = RGBColor(0xFF, 0xFF, 0xFF)

# Tuplas para fondos de celda (la API XML necesita int r,g,b)
AZUL_MED_T  = (0x2E, 0x75, 0xB6)
AZUL_CLAR_T = (0xCA, 0xDC, 0xFC)
BLANCO_T    = (0xFF, 0xFF, 0xFF)
NEGRO     = RGBColor(0x21, 0x21, 0x21)
GRIS      = RGBColor(0x60, 0x60, 0x60)
GRIS_CLAR = RGBColor(0xF4, 0xF4, 0xF4)

W = Inches(13.33)   # LAYOUT_WIDE
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────
def nueva_pres():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank = prs.slide_layouts[6]   # Completamente en blanco
    return prs.slides.add_slide(blank)


def rect(slide, x, y, w, h, fill_rgb, *, alpha=None):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.line.fill.background()
    if fill_rgb:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_rgb
    else:
        shp.fill.background()
    return shp


def txbox(slide, text, x, y, w, h, *,
          size=18, bold=False, color=BLANCO,
          align=PP_ALIGN.LEFT, font="Calibri", wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = size and font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return box


def tabla_slide(slide, headers, rows, x, y, w, h, col_ratios):
    """Tabla sencilla con cabecera azul y filas alternas."""
    cols = len(headers)
    total_rows = 1 + len(rows)
    total_w = Inches(w)
    col_ws = [int(total_w * r / sum(col_ratios)) for r in col_ratios]

    from pptx.util import Emu
    tbl = slide.shapes.add_table(total_rows, cols,
                                  Inches(x), Inches(y),
                                  Inches(w), Inches(h)).table

    # Anchos de columna
    for ci, cw in enumerate(col_ws):
        tbl.columns[ci].width = cw

    def celda(ri, ci, texto, bg=None, txt_color=NEGRO, bold=False, size=11):
        cell = tbl.cell(ri, ci)
        cell.text = texto
        tf = cell.text_frame
        tf.paragraphs[0].font.size  = Pt(size)
        tf.paragraphs[0].font.bold  = bold
        tf.paragraphs[0].font.color.rgb = txt_color
        if bg:
            fill = cell._tc.get_or_add_tcPr()
            solidFill = etree.SubElement(fill, qn("a:solidFill"))
            srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
            srgb.set("val", f"{bg[0]:02X}{bg[1]:02X}{bg[2]:02X}")

    # Cabecera
    for ci, h_text in enumerate(headers):
        celda(0, ci, h_text, bg=AZUL_MED_T, txt_color=BLANCO, bold=True)

    # Filas
    for ri, row in enumerate(rows):
        bg = AZUL_CLAR_T if ri % 2 == 0 else BLANCO_T
        for ci, val in enumerate(row):
            celda(ri + 1, ci, str(val), bg=bg)

    return tbl


def bullet_items(slide, items, x, y, w, h, size=15):
    """Lista de puntos con símbolo •."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"  {item}"
        run.font.size  = Pt(size)
        run.font.color.rgb = NEGRO
        p.space_before = Pt(4)
        # bullet manual (python-pptx no tiene API directa para bullets en textbox)
        pPr = p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "•")
    return box


# ── Slides ────────────────────────────────────────────────────────────────────

def slide_portada(prs):
    sl = blank_slide(prs)
    # Fondo oscuro
    rect(sl, 0, 0, 13.33, 7.5, AZUL_OSC)
    # Franja decorativa derecha
    rect(sl, 11.5, 0, 1.83, 7.5, AZUL_MED)
    # Título
    txbox(sl, "Análisis y Diseño",
          0.6, 2.2, 10.5, 1.2,
          size=48, bold=True, color=BLANCO, align=PP_ALIGN.LEFT, font="Calibri")
    # Subtítulo
    txbox(sl, "Lista de Tareas · DemoCopilot",
          0.6, 3.5, 10.5, 0.7,
          size=24, color=AZUL_CLAR, align=PP_ALIGN.LEFT, font="Calibri")
    # Fecha
    txbox(sl, "Junio 2026",
          0.6, 6.7, 4, 0.4,
          size=12, color=GRIS, align=PP_ALIGN.LEFT, font="Calibri")


def cabecera_slide(sl, titulo):
    """Franja azul superior con título."""
    rect(sl, 0, 0, 13.33, 0.9, AZUL_MED)
    txbox(sl, titulo,
          0.3, 0.05, 12.5, 0.8,
          size=26, bold=True, color=BLANCO, align=PP_ALIGN.LEFT, font="Calibri")
    # Fondo claro resto
    rect(sl, 0, 0.9, 13.33, 6.6, GRIS_CLAR)


def slide_objetivo(prs):
    sl = blank_slide(prs)
    cabecera_slide(sl, "1. Objetivo del proyecto")
    txbox(sl,
          "Aplicación web CRUD de gestión de tareas personales construida como demo didáctica "
          "para el curso de GitHub Copilot.\n\n"
          "Permite crear, consultar, actualizar y eliminar tareas, con soporte de plantillas "
          "reutilizables y tareas repetitivas con generación automática de la siguiente "
          "ocurrencia al completar.\n\n"
          "El objetivo prioritario es la claridad del código sobre la sofisticación arquitectónica.",
          1.0, 1.1, 11.3, 5.8,
          size=16, color=NEGRO, align=PP_ALIGN.LEFT, font="Calibri")


def slide_stack(prs):
    sl = blank_slide(prs)
    cabecera_slide(sl, "2. Stack tecnológico")
    tabla_slide(sl,
        headers=["Tecnología", "Versión", "Rol"],
        rows=[
            ["ASP.NET Core", "8", "API REST / backend (Controllers)"],
            ["Entity Framework Core", "8", "ORM / acceso a datos"],
            ["SQLite", "—", "Base de datos embebida, sin infraestructura externa"],
            ["xUnit + Moq", "—", "Tests unitarios"],
        ],
        x=0.5, y=1.1, w=12.3, h=3.5,
        col_ratios=[2.5, 1, 4])


def slide_arquitectura(prs):
    sl = blank_slide(prs)
    cabecera_slide(sl, "3. Arquitectura de capas")

    tabla_slide(sl,
        headers=["Capa", "Carpeta", "Responsabilidad"],
        rows=[
            ["Modelos de dominio", "Models/",        "Entidades del negocio (TodoItem, PlantillaTarea…)"],
            ["Acceso a datos",     "Data/",           "DbContext y configuración de EF Core"],
            ["Lógica de negocio",  "LogicaNegocio/",  "IXxxLogica + XxxLogica — reglas y acceso a datos"],
            ["Servicios",          "Services/",       "IXxxService + XxxService — orquestación y mapeo"],
            ["Controladores",      "Controllers/",    "Orquestación HTTP, sin lógica de negocio"],
            ["Tests",              "Tests/",          "Proyecto xUnit separado"],
        ],
        x=0.5, y=1.1, w=12.3, h=4.4,
        col_ratios=[2.5, 2, 5])

    txbox(sl, "Flujo:  Controller → [DTO] → Service → [Entidad] → LogicaNegocio → DbContext",
          0.5, 5.8, 12.3, 0.5,
          size=13, bold=True, color=AZUL_MED, font="Courier New")


def slide_modelo(prs):
    sl = blank_slide(prs)
    cabecera_slide(sl, "4. Modelo de datos — TodoItem")
    tabla_slide(sl,
        headers=["Campo", "Tipo", "Descripción"],
        rows=[
            ["Id",               "int",              "Clave primaria, autogenerada"],
            ["Title",            "string",           "Título de la tarea. Requerido"],
            ["IsCompleted",      "bool",             "Estado completada. Por defecto false"],
            ["CreatedAt",        "DateTime",         "Fecha de creación"],
            ["EsRepetitiva",     "bool",             "Si la tarea se repite al completarse"],
            ["Recurrencia",      "TipoRecurrencia?", "Diaria / Semanal / Mensual"],
            ["ProximaFecha",     "DateTime?",        "Siguiente ocurrencia calculada"],
            ["PlantillaId",      "int?",             "FK opcional a PlantillaTarea"],
            ["UsuarioAsignadoId","int?",             "FK opcional a UsuarioAsignado"],
        ],
        x=0.5, y=1.1, w=12.3, h=5.8,
        col_ratios=[3, 2.5, 5])


def slide_entidades(prs):
    sl = blank_slide(prs)
    cabecera_slide(sl, "4. Modelo de datos — Entidades relacionadas")

    txbox(sl, "UsuarioAsignado", 0.5, 1.0, 5.8, 0.5, size=17, bold=True, color=AZUL_MED, font="Calibri")
    tabla_slide(sl,
        headers=["Campo", "Tipo"],
        rows=[
            ["Id",     "int"],
            ["Nombre", "string (máx 100)"],
            ["Email",  "string (único, máx 200)"],
        ],
        x=0.5, y=1.55, w=5.8, h=2.5,
        col_ratios=[2, 3])

    txbox(sl, "PlantillaTarea", 7.0, 1.0, 5.8, 0.5, size=17, bold=True, color=AZUL_MED, font="Calibri")
    tabla_slide(sl,
        headers=["Campo", "Tipo"],
        rows=[
            ["Id",           "int"],
            ["Titulo",       "string"],
            ["EsRepetitiva", "bool"],
            ["Recurrencia",  "TipoRecurrencia?"],
        ],
        x=7.0, y=1.55, w=5.8, h=2.9,
        col_ratios=[2, 3])

    txbox(sl, "TipoRecurrencia (enum):  Diaria  ·  Semanal  ·  Mensual",
          0.5, 4.4, 12.3, 0.5,
          size=15, bold=True, color=NEGRO, font="Calibri")

    txbox(sl,
          "Relación: UsuarioAsignado 1 → N TodoItem  (FK nullable, SetNull al eliminar)\n"
          "Relación: PlantillaTarea 1 → N TodoItem  (FK nullable)",
          0.5, 5.2, 12.3, 1.0,
          size=13, color=GRIS, font="Calibri")


def slide_endpoints_tareas(prs):
    sl = blank_slide(prs)
    cabecera_slide(sl, "5. Endpoints — /api/tareas")
    tabla_slide(sl,
        headers=["Verbo", "Ruta", "Descripción", "OK", "Error"],
        rows=[
            ["GET",    "/api/tareas",               "Listar todas las tareas",                              "200", "—"],
            ["GET",    "/api/tareas/{id}",           "Obtener tarea por ID",                                "200", "404"],
            ["POST",   "/api/tareas",                "Crear nueva tarea",                                   "201", "400"],
            ["PUT",    "/api/tareas/{id}",           "Actualizar título o estado",                          "200", "404/400"],
            ["DELETE", "/api/tareas/{id}",           "Eliminar tarea",                                      "204", "404"],
            ["POST",   "/api/tareas/{id}/completar", "Completar; si repetitiva → genera siguiente ocurrencia", "200", "404"],
        ],
        x=0.5, y=1.1, w=12.3, h=4.5,
        col_ratios=[1.5, 3.5, 5, 1, 1])


def slide_endpoints_otros(prs):
    sl = blank_slide(prs)
    cabecera_slide(sl, "5. Endpoints — Usuarios y Plantillas")

    txbox(sl, "/api/usuariosasignados", 0.5, 1.0, 12.3, 0.4, size=15, bold=True, color=AZUL_MED)
    tabla_slide(sl,
        headers=["Verbo", "Ruta", "Descripción"],
        rows=[
            ["GET / POST",   "/api/usuariosasignados",      "Listar / Crear usuario"],
            ["GET / PUT",    "/api/usuariosasignados/{id}", "Obtener / Actualizar usuario"],
            ["DELETE",       "/api/usuariosasignados/{id}", "Eliminar (tareas quedan sin usuario)"],
        ],
        x=0.5, y=1.45, w=12.3, h=2.0,
        col_ratios=[2, 3.5, 5])

    txbox(sl, "/api/plantillas", 0.5, 3.65, 12.3, 0.4, size=15, bold=True, color=AZUL_MED)
    tabla_slide(sl,
        headers=["Verbo", "Ruta", "Descripción"],
        rows=[
            ["GET / POST",   "/api/plantillas",                  "Listar / Crear plantilla"],
            ["GET / PUT",    "/api/plantillas/{id}",             "Obtener / Actualizar plantilla"],
            ["DELETE",       "/api/plantillas/{id}",             "Eliminar plantilla"],
            ["POST",         "/api/plantillas/{id}/instanciar",  "Crear tarea a partir de la plantilla"],
        ],
        x=0.5, y=4.1, w=12.3, h=2.6,
        col_ratios=[2, 3.5, 5])


def slide_decisiones(prs):
    sl = blank_slide(prs)
    cabecera_slide(sl, "7. Decisiones de diseño")
    bullet_items(sl, [
        "SQLite sobre SQL Server: embebida, sin infraestructura. Ideal para demo local.",
        "Controllers sobre Minimal API: más explícitos y legibles en pantalla.",
        "Inyección de dependencias: desacopla lógica y facilita tests con mocks.",
        "POST /completar en lugar de PUT: la acción tiene efecto secundario (nueva ocurrencia).",
        "Plantilla como entidad independiente: permite gestionar plantillas sin acoplarlas a tareas.",
        "Recurrencia calculada en el servicio: lógica en TodoService.Completar(), no en el controlador.",
        "Sin paginación ni autenticación: fuera del alcance de la demo.",
        "Código en castellano: coherencia con el público hispanohablante del curso.",
    ], x=0.8, y=1.1, w=11.7, h=6.0, size=15)


def slide_pendientes(prs):
    sl = blank_slide(prs)
    cabecera_slide(sl, "8. Pendientes / Preguntas abiertas")
    tabla_slide(sl,
        headers=["Tema", "Detalle"],
        rows=[
            ["Frontend",                    "¿Razor Pages o React + Vite? Pendiente de decisión."],
            ["Paginación",                  "No incluida. Se añadiría skip/take si crece el volumen."],
            ["Autenticación",               "No incluida. Extensión futura con JWT o cookies."],
            ["Filtros y búsqueda",          "Posible: GET /api/tareas?completada=true."],
            ["Soft delete",                 "Actualmente hard delete. Campo DeletedAt si se requiere historial."],
            ["Job de recurrencia",          "Solo se genera al llamar a POST /completar. BackgroundService si se necesita automático."],
            ["Recurrencia personalizada",   "Enum cubre Diaria/Semanal/Mensual. Extensión pendiente si se requieren patrones complejos."],
            ["Herencia plantilla → tareas", "Modificar plantilla no actualiza tareas existentes. Comportamiento a definir."],
        ],
        x=0.5, y=1.1, w=12.3, h=5.8,
        col_ratios=[3, 7])


def slide_cierre(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, AZUL_OSC)
    rect(sl, 0, 0, 1.5, 7.5, AZUL_MED)
    txbox(sl, "DemoCopilot",
          2.0, 2.5, 9, 1.0,
          size=40, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, font="Calibri")
    txbox(sl, "github.com/tu-usuario/DemoCopilot",
          2.0, 3.7, 9, 0.6,
          size=16, color=AZUL_CLAR, align=PP_ALIGN.CENTER, font="Calibri")
    txbox(sl, "Construido con  ASP.NET Core 8  ·  EF Core  ·  SQLite  ·  GitHub Copilot",
          2.0, 6.5, 9, 0.5,
          size=12, color=GRIS, align=PP_ALIGN.CENTER, font="Calibri")


# ── main ──────────────────────────────────────────────────────────────────────
def build():
    prs = nueva_pres()

    slide_portada(prs)
    slide_objetivo(prs)
    slide_stack(prs)
    slide_arquitectura(prs)
    slide_modelo(prs)
    slide_entidades(prs)
    slide_endpoints_tareas(prs)
    slide_endpoints_otros(prs)
    slide_decisiones(prs)
    slide_pendientes(prs)
    slide_cierre(prs)

    prs.save(OUTPUT)
    print(f"Presentación generada: {OUTPUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
